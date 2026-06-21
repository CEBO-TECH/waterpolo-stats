#!/usr/bin/env python3
"""Generate the Cap Track marketing deck (PZPW) as a .pptx.

Dark, on-brand (Cap Track theme) Polish slides with stylized in-app mockups
drawn as vector shapes. Run: python3 marketing/generate_deck.py
Output: marketing/Cap-Track-PZPW.pptx
"""

import os

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand palette (matches frontend/app/globals.css) ──
BG = RGBColor(0x0B, 0x0F, 0x14)
CARD = RGBColor(0x13, 0x1A, 0x24)
CARD2 = RGBColor(0x1A, 0x23, 0x32)
BORDER = RGBColor(0x1E, 0x2D, 0x3D)
ACCENT = RGBColor(0x4C, 0xC9, 0xF0)
ACCENT_DK = RGBColor(0x14, 0x53, 0x66)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x97, 0xA5)
FAINT = RGBColor(0x5F, 0x6B, 0x7A)
GREEN = RGBColor(0x51, 0xCF, 0x66)
RED = RGBColor(0xFF, 0x6B, 0x6B)
YELLOW = RGBColor(0xFF, 0xD4, 0x3B)
DARK = RGBColor(0x00, 0x10, 0x18)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, radius=None):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def oval(s, l, t, w, h, fill, line=None, lw=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def tb(s, l, t, w, h, paras, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", align)
        if "after" in para:
            p.space_after = Pt(para["after"])
        if "before" in para:
            p.space_before = Pt(para["before"])
        p.line_spacing = para.get("line", 1.0)
        for run in para["runs"]:
            r = p.add_run()
            r.text = run[0]
            r.font.size = Pt(run[1])
            r.font.color.rgb = run[2]
            r.font.bold = run[3] if len(run) > 3 else False
            r.font.name = FONT
    return box


def label_in(s, sp_l, sp_t, sp_w, sp_h, text, size, color, bold=True, align=PP_ALIGN.CENTER):
    return tb(s, sp_l, sp_t, sp_w, sp_h, [{"runs": [(text, size, color, bold)], "align": align}],
              anchor=MSO_ANCHOR.MIDDLE)


def scaffold(s, kicker, title, page):
    rect(s, 0.62, 0.66, 0.07, 0.62, fill=ACCENT)
    tb(s, 0.85, 0.6, 11.5, 0.4, [{"runs": [(kicker.upper(), 13, ACCENT, True)]}])
    tb(s, 0.83, 0.95, 11.8, 1.0, [{"runs": [(title, 30, TEXT, True)], "line": 1.0}])
    # footer
    rect(s, 0.62, 6.92, 12.1, 0.012, fill=BORDER)
    tb(s, 0.62, 6.98, 6, 0.34,
       [{"runs": [("Cap Track", 11, ACCENT, True), ("   —  inteligentne statystyki waterpolo", 11, FAINT, False)]}])
    tb(s, 8.0, 6.98, 4.7, 0.34,
       [{"runs": [(f"Powered by CEBO.TECH   ·   {page:02d}", 11, FAINT, False)], "align": PP_ALIGN.RIGHT}])


def bullets(s, l, t, w, items, size=16, gap=10, line=1.05):
    paras = []
    for it in items:
        head = it[0] if isinstance(it, tuple) else it
        sub = it[1] if isinstance(it, tuple) and len(it) > 1 else None
        runs = [("▸  ", size, ACCENT, True), (head, size, TEXT, True)]
        if sub:
            runs.append(("  " + sub, size, MUTED, False))
        paras.append({"runs": runs, "after": gap, "line": line})
    return tb(s, l, t, w, 2.9, paras)


def chip(s, l, t, w, h, text, fill, txtcolor, size=12, bold=True, radius=0.5, line=None):
    rect(s, l, t, w, h, fill=fill, radius=radius, line=line)
    label_in(s, l, t, w, h, text, size, txtcolor, bold=bold)


def logo(s, cx, cy, r=0.42, wordmark=True, size=14):
    oval(s, cx - r, cy - r, 2 * r, 2 * r, fill=ACCENT)
    label_in(s, cx - r, cy - r, 2 * r, 2 * r, "13", size, DARK, bold=True)
    if wordmark:
        tb(s, cx + r + 0.12, cy - r - 0.05, 4, 2 * r + 0.1,
           [{"runs": [("Cap Track", size + 8, TEXT, True)]}], anchor=MSO_ANCHOR.MIDDLE)


# ── Mockups ──

def m_frame(s, l, t, w, h, title):
    rect(s, l, t, w, h, fill=CARD, line=BORDER, lw=1.25, radius=0.04)
    rect(s, l, t, w, 0.46, fill=CARD2, radius=0.06)
    label_in(s, l + 0.2, t, w - 0.4, 0.46, title, 11, MUTED, bold=True, align=PP_ALIGN.LEFT)
    for i, c in enumerate((RED, YELLOW, GREEN)):
        oval(s, l + w - 0.28 - i * 0.22, t + 0.17, 0.12, 0.12, fill=c)


def m_scorekeeper(s, l, t, w, h):
    m_frame(s, l, t, w, h, "Asystent meczowy — na żywo")
    iy = t + 0.62
    # left lineup column
    lw = 1.7
    label_in(s, l + 0.18, iy, lw, 0.3, "WODA", 10, ACCENT, align=PP_ALIGN.LEFT)
    for i in range(4):
        cy = iy + 0.38 + i * 0.5
        rect(s, l + 0.18, cy, lw, 0.42, fill=CARD2, line=BORDER, radius=0.18)
        oval(s, l + 0.28, cy + 0.07, 0.28, 0.28, fill=ACCENT_DK, line=ACCENT, lw=1)
        label_in(s, l + 0.28, cy + 0.07, 0.28, 0.28, str([12, 7, 4, 9][i]), 10, ACCENT)
        rect(s, l + 0.72, cy + 0.16, 0.7, 0.12, fill=BORDER, radius=0.5)
    # right action grid
    gx = l + 2.1
    gw = w - 2.35
    cols = 4
    labels = ["GOL akcja", "GOL kontra", "Asysta", "Obrona", "Strata", "Niecelny",
              "Wykl.", "Karny", "Przejęcie", "Blok", "Centra", "5m"]
    cw = (gw - (cols - 1) * 0.12) / cols
    for i, lab in enumerate(labels):
        r, c = divmod(i, cols)
        bx = gx + c * (cw + 0.12)
        by = iy + 0.05 + r * 0.52
        fill = ACCENT if i in (0, 1) else CARD2
        tcol = DARK if i in (0, 1) else TEXT
        chip(s, bx, by, cw, 0.42, lab, fill, tcol, size=9.5, line=None if i in (0, 1) else BORDER)
    # voice button across bottom
    vy = iy + 0.05 + 3 * 0.52 + 0.12
    rect(s, gx, vy, gw, 0.5, fill=RGBColor(0x10, 0x2A, 0x33), line=ACCENT, lw=1.25, radius=0.12)
    label_in(s, gx, vy, gw, 0.5, "🎤  Komenda głosem", 12, ACCENT)


def m_voice(s, l, t, w, h):
    m_frame(s, l, t, w, h, "Sterowanie głosem")
    iy = t + 0.66
    # listening mic
    rect(s, l + 0.25, iy, w - 0.5, 0.6, fill=RGBColor(0x2A, 0x12, 0x14), line=RED, lw=1.25, radius=0.14)
    label_in(s, l + 0.25, iy, w - 0.5, 0.6, "🎤  Słucham…", 14, RED)
    # confirm chip
    cy = iy + 0.78
    ch = h - (iy - t) - 0.95
    rect(s, l + 0.25, cy, w - 0.5, ch, fill=CARD2, line=ACCENT, lw=1.25, radius=0.08)
    oval(s, l + 0.45, cy + 0.22, 0.5, 0.5, fill=ACCENT)
    label_in(s, l + 0.45, cy + 0.22, 0.5, 0.5, "12", 15, DARK)
    tb(s, l + 1.08, cy + 0.2, w - 1.9, 0.55,
       [{"runs": [("Gol z kontrataku", 16, TEXT, True)]}], anchor=MSO_ANCHOR.MIDDLE)
    chip(s, l + w - 0.95, cy + 0.28, 0.55, 0.34, "AI", RGBColor(0x10, 0x2A, 0x33), ACCENT, size=11, radius=0.5, line=ACCENT)
    tb(s, l + 0.45, cy + 0.78, w - 0.9, 0.3,
       [{"runs": [("„dwunastka trafiła z kontry”", 11.5, MUTED, False)]}])
    by = cy + ch - 0.62
    chip(s, l + 0.45, by, (w - 1.0) / 2 - 0.06, 0.4, "✓ Zapisz", ACCENT, DARK, size=12, radius=0.2)
    chip(s, l + 0.45 + (w - 1.0) / 2 + 0.06, by, (w - 1.0) / 2 - 0.06, 0.4, "✕ Anuluj", CARD, TEXT, size=12, radius=0.2, line=BORDER)
    rect(s, l + 0.45, cy + ch - 0.14, w - 0.9, 0.06, fill=BORDER, radius=0.5)
    rect(s, l + 0.45, cy + ch - 0.14, (w - 0.9) * 0.62, 0.06, fill=ACCENT, radius=0.5)


def m_bars(s, l, t, w, h, values, title, color=ACCENT):
    m_frame(s, l, t, w, h, title)
    base = t + h - 0.5
    top = t + 0.75
    avail_h = base - top
    n = len(values)
    bw = (w - 0.8) / n
    mx = max(values)
    for i, v in enumerate(values):
        bh = avail_h * (v / mx)
        bx = l + 0.4 + i * bw + bw * 0.18
        rect(s, bx, base - bh, bw * 0.64, bh, fill=color if i != values.index(mx) else ACCENT, radius=0.12)
    rect(s, l + 0.35, base, w - 0.7, 0.02, fill=BORDER)


def m_dashboard(s, l, t, w, h):
    m_frame(s, l, t, w, h, "Pulpit sezonu")
    iy = t + 0.66
    kpis = [("MECZE", "18"), ("BILANS", "12-2-4"), ("GOLE +", "214"), ("GOLE −", "151")]
    kw = (w - 0.5 - 3 * 0.16) / 4
    for i, (lab, val) in enumerate(kpis):
        kx = l + 0.25 + i * (kw + 0.16)
        rect(s, kx, iy, kw, 1.0, fill=CARD2, radius=0.1)
        tb(s, kx, iy + 0.14, kw, 0.3, [{"runs": [(lab, 9, MUTED, True)], "align": PP_ALIGN.CENTER}])
        tb(s, kx, iy + 0.42, kw, 0.5, [{"runs": [(val, 22, ACCENT, True)], "align": PP_ALIGN.CENTER}])
    # mini bar row
    by = iy + 1.3
    tb(s, l + 0.25, by, 4, 0.3, [{"runs": [("Forma — ostatnie 6", 10, MUTED, True)]}])
    res = [(GREEN, "W"), (GREEN, "W"), (RED, "P"), (GREEN, "W"), (YELLOW, "R"), (GREEN, "W")]
    for i, (c, t2) in enumerate(res):
        cx = l + 0.25 + i * 0.55
        rect(s, cx, by + 0.34, 0.42, 0.42, fill=c, radius=0.2)
        label_in(s, cx, by + 0.34, 0.42, 0.42, t2, 12, DARK)
    # top scorers
    tb(s, l + 4.2, by, 4, 0.3, [{"runs": [("Najlepsi strzelcy", 10, MUTED, True)]}])
    for i, (nm, g) in enumerate([("#12 Kowalski", 41), ("#7 Nowak", 33), ("#9 Wiśniewski", 28)]):
        ry = by + 0.34 + i * 0.42
        rect(s, l + 4.2, ry, w - 4.6, 0.34, fill=CARD2, radius=0.16)
        tb(s, l + 4.4, ry, w - 5.6, 0.34, [{"runs": [(nm, 11, TEXT, False)]}], anchor=MSO_ANCHOR.MIDDLE)
        tb(s, l + w - 1.3, ry, 0.9, 0.34, [{"runs": [(str(g), 11, ACCENT, True)], "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)


def m_mvp(s, l, t, w, h):
    m_frame(s, l, t, w, h, "MVP po meczu")
    iy = t + 0.66
    rect(s, l + 0.25, iy, w - 0.5, 1.25, fill=RGBColor(0x10, 0x2A, 0x33), line=ACCENT, lw=1.25, radius=0.08)
    oval(s, l + 0.5, iy + 0.32, 0.62, 0.62, fill=ACCENT)
    label_in(s, l + 0.5, iy + 0.32, 0.62, 0.62, "★", 22, DARK)
    tb(s, l + 1.3, iy + 0.28, w - 2, 0.5, [{"runs": [("#12 Jan Kowalski", 18, TEXT, True)]}])
    tb(s, l + 1.3, iy + 0.72, w - 2, 0.4, [{"runs": [("3 gole · 2 asysty · 4 przejęcia", 12, MUTED, False)]}])
    chip(s, l + w - 1.5, iy + 0.45, 1.0, 0.5, "98 pkt", ACCENT, DARK, size=15, radius=0.2)
    for i, (nm, p) in enumerate([("#7 Nowak", 81), ("#9 Wiśniewski", 74), ("#4 Lewandowski", 63)]):
        ry = iy + 1.45 + i * 0.5
        rect(s, l + 0.25, ry, w - 0.5, 0.42, fill=CARD2, radius=0.14)
        tb(s, l + 0.45, ry, w - 2, 0.42, [{"runs": [(f"{i + 2}.  {nm}", 12, TEXT, False)]}], anchor=MSO_ANCHOR.MIDDLE)
        tb(s, l + w - 1.4, ry, 1.0, 0.42, [{"runs": [(f"{p} pkt", 12, MUTED, True)], "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)


SCREENS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "screens")


def shot(s, l, t, w, h, name, title, fallback=None):
    """Place a real app screenshot inside a window frame, preserving aspect.

    Falls back to the stylized vector mockup if the PNG is missing, so the deck
    still builds without the screenshots present.
    """
    path = os.path.join(SCREENS, name + ".png")
    if not os.path.exists(path):
        if fallback:
            fallback(s, l, t, w, h)
        return
    m_frame(s, l, t, w, h, title)
    il, it = l + 0.12, t + 0.56
    iw, ih = w - 0.24, h - 0.70
    pw, ph = PILImage.open(path).size
    ratio = ph / pw
    if iw * ratio <= ih:
        dw, dh = iw, iw * ratio
    else:
        dw, dh = ih / ratio, ih
    dx, dy = il + (iw - dw) / 2, it + (ih - dh) / 2
    s.shapes.add_picture(path, Inches(dx), Inches(dy), Inches(dw), Inches(dh))
    rect(s, dx, dy, dw, dh, fill=None, line=BORDER, lw=0.75, radius=0.015)


def feature_card(s, l, t, w, h, title, desc):
    rect(s, l, t, w, h, fill=CARD, line=BORDER, lw=1, radius=0.08)
    rect(s, l + 0.22, t + 0.24, 0.34, 0.34, fill=ACCENT, radius=0.3)
    tb(s, l + 0.7, t + 0.2, w - 0.9, 0.45, [{"runs": [(title, 14, TEXT, True)]}], anchor=MSO_ANCHOR.MIDDLE)
    tb(s, l + 0.24, t + 0.74, w - 0.48, h - 0.9, [{"runs": [(desc, 11, MUTED, False)], "line": 1.05}])


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, W, 0.14, fill=ACCENT)
logo(s, 6.67, 2.55, r=0.55, wordmark=False, size=22)
tb(s, 0, 3.25, W, 0.9, [{"runs": [("Cap Track", 52, TEXT, True)], "align": PP_ALIGN.CENTER}])
tb(s, 0, 4.2, W, 0.6, [{"runs": [("Inteligentna analiza statystyk piłki wodnej", 20, ACCENT, False)], "align": PP_ALIGN.CENTER}])
tb(s, 0, 4.85, W, 0.5,
   [{"runs": [("Rejestracja na żywo · statystyki · wideo · AI · aplikacja mobilna", 14, MUTED, False)], "align": PP_ALIGN.CENTER}])
chip(s, 4.67, 5.75, 4.0, 0.55, "Prezentacja dla Polskiego Związku Piłki Wodnej", CARD, TEXT, size=13, radius=0.3, line=BORDER)
tb(s, 0, 7.0, W, 0.4, [{"runs": [("Powered by CEBO.TECH", 12, FAINT, False)], "align": PP_ALIGN.CENTER}])

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ════════════════════════════════════════════════════════════════
s = slide()
scaffold(s, "Wyzwanie", "Statystyki waterpolo robi się „na kartce”", 2)
bullets(s, 0.85, 2.1, 6.0, [
    ("Zeszyt, Excel, pamięć trenera", "— dane niespójne i łatwe do zgubienia."),
    ("Brak narzędzia dedykowanego waterpolo", "— aplikacje z innych sportów nie pasują."),
    ("Analiza po meczu zajmuje godziny", "— ręczne przepisywanie i liczenie."),
    ("Brak obiektywnych danych", "— decyzje kadrowe „na oko”."),
    ("Wideo żyje osobno od statystyk", "— trudno wrócić do konkretnej akcji."),
], size=15, gap=14)
m_frame(s, 7.4, 2.1, 5.3, 4.4, "Dziś: notatnik trenera")
for i in range(7):
    ry = 2.95 + i * 0.5
    rect(s, 7.75, ry, 4.6, 0.34, fill=CARD2, radius=0.1)
    rect(s, 7.9, ry + 0.11, 2.2 + (i % 3) * 0.7, 0.12, fill=BORDER, radius=0.5)
    label_in(s, 11.7, ry, 0.5, 0.34, "?", 13, FAINT)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Rozwiązanie
# ════════════════════════════════════════════════════════════════
s = slide()
scaffold(s, "Rozwiązanie", "Cap Track — jedno miejsce na cały mecz", 3)
tb(s, 0.85, 2.05, 5.9, 2.2, [
    {"runs": [("Aplikacja zaprojektowana ", 16, TEXT, False), ("pod waterpolo", 16, ACCENT, True),
              (" i pod ", 16, TEXT, False), ("iPada", 16, ACCENT, True),
              (". Trener jednym dotykiem (lub głosem) zapisuje zdarzenie, a system od razu liczy statystyki, czas gry i podpina wideo.", 16, TEXT, False)],
     "line": 1.25, "after": 14},
])
bullets(s, 0.85, 3.6, 5.9, [
    ("Zbieranie danych w czasie rzeczywistym", ""),
    ("44 typy zdarzeń meczowych", "— atak, obrona, przewagi, karne."),
    ("Działa offline na hali", "— bez ryzyka utraty danych."),
    ("Web + iOS + Android", "— jeden kod, każde urządzenie."),
], size=15, gap=12)
shot(s, 7.15, 1.95, 5.55, 4.6, "asystent", "Asystent meczowy — na żywo", fallback=m_scorekeeper)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Funkcje (grid)
# ════════════════════════════════════════════════════════════════
s = slide()
scaffold(s, "Przegląd", "Kompletny zestaw narzędzi dla sztabu", 4)
feats = [
    ("Asystent na żywo", "Szybkie zapisywanie zdarzeń meczowych jednym dotykiem."),
    ("Sterowanie głosem AI", "„Numer 12 gol z kontrataku” — agent dopisuje zdarzenie."),
    ("Tryb offline", "Komendy nie giną przy zerwaniu sieci na hali."),
    ("Wideo + YouTube", "Każde zdarzenie linkuje do klipu z właściwej sekundy."),
    ("Statystyki zaawansowane", "Wykresy, multi-mecz, indeks ataku i obrony."),
    ("Czas gry (woda/ławka)", "Automatyczne liczenie czasu w wodzie per zawodnik."),
    ("MVP po meczu", "Obiektywny ranking wkładu zawodników."),
    ("Pulpit sezonu", "KPI, forma, najlepsi strzelcy i asystenci."),
    ("Role i kluby", "Właściciel, trener, zawodnik · wiele klubów."),
]
gw = (12.1 - 2 * 0.3) / 3
gh = 1.35
for i, (ti, de) in enumerate(feats):
    r, c = divmod(i, 3)
    feature_card(s, 0.62 + c * (gw + 0.3), 2.05 + r * (gh + 0.25), gw, gh, ti, de)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Asystent na żywo
# ════════════════════════════════════════════════════════════════
s = slide()
scaffold(s, "Funkcja · 01", "Asystent meczowy — klikasz, system liczy", 5)
shot(s, 0.62, 2.0, 6.4, 4.5, "asystent", "Asystent meczowy — na żywo", fallback=m_scorekeeper)
bullets(s, 7.4, 2.2, 5.3, [
    ("Duże przyciski pod iPada", "— wygodne tempo meczu."),
    ("Podział WODA / ŁAWKA", "— ze strzałkami wejście/zejście."),
    ("Atak pozycyjny ↔ przewaga", "— jeden przełącznik."),
    ("Cofnij ostatnią akcję", "— błąd poprawisz w sekundę."),
    ("Wynik liczony automatycznie", "— po kwartach i łącznie."),
], size=15, gap=14)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Voice AI
# ════════════════════════════════════════════════════════════════
s = slide()
scaffold(s, "Funkcja · 02", "Sterowanie głosem — ręce wolne na meczu", 6)
m_voice(s, 0.62, 2.0, 5.4, 4.5)
tb(s, 6.5, 2.15, 6.2, 1.0, [
    {"runs": [("Trener patrzy na wodę, nie na ekran. ", 16, TEXT, False),
              ("Mówisz komendę — agent ją rozumie i dopisuje zdarzenie.", 16, ACCENT, True)], "line": 1.25, "after": 12},
])
bullets(s, 6.5, 3.5, 6.2, [
    ("Rozpoznawanie mowy po polsku", "— wbudowane w przeglądarkę."),
    ("Parser działa natychmiast i offline", "— bez kosztów, w pełni przewidywalnie."),
    ("Warstwa AI (Claude) dla swobodnych zdań", "— „dwunastka trafiła z kontry”."),
    ("Potwierdzenie przed zapisem", "— chroni dane przed przesłyszeniem."),
], size=15, gap=13)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Offline
# ════════════════════════════════════════════════════════════════
s = slide()
scaffold(s, "Funkcja · 03", "Działa offline — nic nie ginie", 7)
bullets(s, 0.85, 2.2, 6.0, [
    ("Hale często mają słaby zasięg", "— Cap Track to przewidział."),
    ("Każdy zapis trafia do kolejki", "— gdy sieć wróci, wysyła się sam."),
    ("Optymistyczny interfejs", "— widzisz akcję od razu, bez czekania."),
    ("Zero utraconych danych", "— spokój podczas całego meczu."),
], size=16, gap=16)
# offline flow mockup
m_frame(s, 7.2, 2.05, 5.5, 4.3, "Kolejka offline")
steps = [("Zapis zdarzenia", ACCENT, "✓"), ("Brak sieci — do kolejki", YELLOW, "…"),
         ("Sieć wraca", ACCENT, "↻"), ("Auto-synchronizacja", GREEN, "✓")]
for i, (txt, c, ic) in enumerate(steps):
    ry = 2.75 + i * 0.85
    rect(s, 7.55, ry, 4.8, 0.6, fill=CARD2, line=BORDER, radius=0.12)
    oval(s, 7.7, ry + 0.13, 0.34, 0.34, fill=c)
    label_in(s, 7.7, ry + 0.13, 0.34, 0.34, ic, 12, DARK)
    tb(s, 8.2, ry, 4.0, 0.6, [{"runs": [(txt, 13, TEXT, False)]}], anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        label_in(s, 9.7, ry + 0.6, 0.5, 0.25, "↓", 12, FAINT)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Wideo
# ════════════════════════════════════════════════════════════════
s = slide()
scaffold(s, "Funkcja · 04", "Wideo zsynchronizowane z każdą akcją", 8)
bullets(s, 0.85, 2.2, 5.6, [
    ("Podpinasz transmisję z YouTube", "— ustawiasz start = teraz."),
    ("Każde zdarzenie ma znacznik czasu", "— dokładnie w sekundzie akcji."),
    ("Jeden klik → klip od właściwego momentu", "— z cofnięciem o kilka sekund."),
    ("Analiza wideo bez przewijania", "— ogromna oszczędność czasu."),
], size=16, gap=15)
# video mockup
m_frame(s, 6.9, 2.05, 5.8, 4.4, "Powtórka akcji")
rect(s, 7.2, 2.7, 5.2, 2.6, fill=DARK, line=BORDER, radius=0.04)
oval(s, 9.45, 3.7, 0.7, 0.7, fill=ACCENT)
label_in(s, 9.45, 3.68, 0.7, 0.7, "▶", 18, DARK)
rect(s, 7.4, 5.55, 4.8, 0.12, fill=BORDER, radius=0.5)
rect(s, 7.4, 5.55, 1.9, 0.12, fill=ACCENT, radius=0.5)
oval(s, 9.2, 5.5, 0.22, 0.22, fill=ACCENT)
tb(s, 7.2, 5.85, 5.2, 0.4, [{"runs": [("#12 Gol z kontrataku · 2. kwarta · 14:32", 11, MUTED, False)], "align": PP_ALIGN.CENTER}])

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Statystyki
# ════════════════════════════════════════════════════════════════
s = slide()
scaffold(s, "Funkcja · 05", "Statystyki, które realnie pomagają trenować", 9)
shot(s, 0.62, 2.0, 6.2, 4.5, "stats", "Statystyki — filtry, tabela, wykresy",
     fallback=lambda s, l, t, w, h: m_bars(s, l, t, w, h, [6, 9, 5, 8, 7, 10, 6, 9], "Skuteczność ataku wg kwarty"))
bullets(s, 7.2, 2.2, 5.5, [
    ("Lekkie wykresy w aplikacji", "— bez eksportu do Excela."),
    ("Analiza wielu meczów naraz", "— trendy formy zespołu."),
    ("Indeks ataku i obrony", "— jedna liczba opisująca grę."),
    ("Rozkład wg kwart", "— gdzie tracimy mecze."),
    ("Porównania zawodników i drużyn", "— obiektywnie, na danych."),
], size=15, gap=14)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Dashboard
# ════════════════════════════════════════════════════════════════
s = slide()
scaffold(s, "Funkcja · 06", "Pulpit — cały sezon na jednym ekranie", 10)
shot(s, 0.62, 2.0, 7.0, 4.5, "pulpit", "Pulpit sezonu", fallback=m_dashboard)
bullets(s, 7.9, 2.3, 4.8, [
    ("Najważniejsze KPI sezonu", "— mecze, bilans, bramki."),
    ("Forma w skrócie", "— ostatnie wyniki W/R/P."),
    ("Najlepsi strzelcy i asystenci", ""),
    ("Szybki powrót do meczu", "— „Wznów” lub „Nowy mecz”."),
], size=15, gap=15)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — MVP + czas gry
# ════════════════════════════════════════════════════════════════
s = slide()
scaffold(s, "Funkcja · 07", "MVP i czas gry — docenianie zawodników", 11)
shot(s, 0.62, 2.0, 6.3, 4.5, "mvp", "MVP po meczu", fallback=m_mvp)
bullets(s, 7.3, 2.3, 5.4, [
    ("Obiektywny ranking wkładu", "— ważony scoring zdarzeń."),
    ("Sugerowany MVP po meczu", "— trener zatwierdza lub zmienia."),
    ("Czas w wodzie per zawodnik", "— liczony automatycznie ze zmian."),
    ("Sprawiedliwy podział minut", "— dane zamiast wrażeń."),
    ("Motywacja dla młodzieży", "— widoczne docenienie pracy."),
], size=15, gap=14)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Role / kluby / konta
# ════════════════════════════════════════════════════════════════
s = slide()
scaffold(s, "Funkcja · 08", "Role, kluby i konta zawodników", 12)
roles = [
    ("Właściciel", "Zarządza klubem, użytkownikami i ustawieniami."),
    ("Trener", "Prowadzi mecze, składy i statystyki zespołu."),
    ("Zawodnik", "Widzi własne statystyki, mecze i czas gry."),
]
cw = (12.1 - 2 * 0.35) / 3
for i, (ti, de) in enumerate(roles):
    cx = 0.62 + i * (cw + 0.35)
    rect(s, cx, 2.1, cw, 1.9, fill=CARD, line=BORDER, lw=1, radius=0.06)
    oval(s, cx + cw / 2 - 0.35, 2.4, 0.7, 0.7, fill=ACCENT_DK, line=ACCENT, lw=1.25)
    label_in(s, cx + cw / 2 - 0.35, 2.4, 0.7, 0.7, ["★", "✦", "●"][i], 18, ACCENT)
    tb(s, cx, 3.25, cw, 0.4, [{"runs": [(ti, 17, TEXT, True)], "align": PP_ALIGN.CENTER}])
    tb(s, cx + 0.3, 3.72, cw - 0.6, 0.9, [{"runs": [(de, 12, MUTED, False)], "line": 1.1, "align": PP_ALIGN.CENTER}])
bullets(s, 0.85, 4.35, 11.6, [
    ("Wiele klubów na jednym koncie", "— przełączasz się jednym kliknięciem."),
    ("Zaproszenia e-mail z rolą", "— szybkie dodawanie sztabu i zawodników."),
    ("Zawodnik = własny widok", "— bez dostępu do paneli zarządzania."),
], size=15, gap=12)

# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Mobile + tech
# ════════════════════════════════════════════════════════════════
s = slide()
scaffold(s, "Technologia", "Aplikacja mobilna i nowoczesny fundament", 13)
# phone mockup
m_frame(s, 0.9, 2.0, 2.6, 4.5, "iPad / iPhone")
rect(s, 1.2, 2.7, 2.0, 3.4, fill=DARK, line=ACCENT, lw=1, radius=0.08)
logo(s, 2.2, 3.45, r=0.3, wordmark=False, size=12)
for i in range(4):
    rect(s, 1.4, 4.1 + i * 0.45, 1.6, 0.32, fill=CARD2 if i else ACCENT, radius=0.2)
bullets(s, 4.1, 2.2, 8.4, [
    ("Web, iOS i Android z jednego kodu", "— natywne aplikacje (Capacitor)."),
    ("Zaprojektowane „iPad-first”", "— główne narzędzie pracy na hali."),
    ("Backend w architekturze heksagonalnej", "— czysty, testowany, łatwy do rozwoju."),
    ("Własny hosting (Coabify/Docker)", "— pełna kontrola nad wdrożeniem."),
    ("Dane na serwerach w UE", "— zgodność z RODO."),
    ("Bezpieczne logowanie i role", "— dostęp tylko dla uprawnionych."),
], size=15, gap=13)

# ════════════════════════════════════════════════════════════════
# SLIDE 14 — Współpraca / kontakt
# ════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, W, 0.14, fill=ACCENT)
tb(s, 0.85, 0.9, 11.6, 0.4, [{"runs": [("WIZJA WSPÓŁPRACY", 14, ACCENT, True)]}])
tb(s, 0.83, 1.3, 11.6, 0.9, [{"runs": [("Ujednolićmy statystyki polskiego waterpolo", 30, TEXT, True)], "line": 1.0}])
cards = [
    ("Kluby i akademie", "Wspólny standard zbierania danych w całym kraju."),
    ("Kadry narodowe", "Obiektywna analiza zawodników i przeciwników."),
    ("Rozwój młodzieży", "Mierzalny postęp i sprawiedliwy czas gry."),
]
cw = (12.1 - 2 * 0.35) / 3
for i, (ti, de) in enumerate(cards):
    cx = 0.62 + i * (cw + 0.35)
    rect(s, cx, 2.5, cw, 1.7, fill=CARD, line=BORDER, radius=0.06)
    rect(s, cx + 0.3, 2.8, 0.45, 0.45, fill=ACCENT, radius=0.3)
    tb(s, cx + 0.3, 3.4, cw - 0.6, 0.4, [{"runs": [(ti, 15, TEXT, True)]}])
    tb(s, cx + 0.3, 3.78, cw - 0.6, 0.8, [{"runs": [(de, 11.5, MUTED, False)], "line": 1.1}])
# contact bar
rect(s, 0.62, 4.7, 12.1, 1.55, fill=CARD, line=ACCENT, lw=1.25, radius=0.05)
logo(s, 1.6, 5.47, r=0.45, wordmark=False, size=16)
tb(s, 2.4, 4.95, 7, 0.5, [{"runs": [("Cap Track", 24, TEXT, True)]}])
tb(s, 2.42, 5.55, 8, 0.5, [{"runs": [("Zróbmy pilotaż w wybranych klubach — pokażemy pełnię możliwości na żywo.", 13, MUTED, False)]}])
tb(s, 8.4, 4.95, 4.0, 1.2,
   [{"runs": [("Kontakt", 12, ACCENT, True)], "after": 4, "align": PP_ALIGN.RIGHT},
    {"runs": [("patrykcebo11@gmail.com", 14, TEXT, True)], "after": 2, "align": PP_ALIGN.RIGHT},
    {"runs": [("CEBO.TECH", 12, MUTED, False)], "align": PP_ALIGN.RIGHT}])
tb(s, 0, 6.95, W, 0.4, [{"runs": [("Powered by CEBO.TECH", 11, FAINT, False)], "align": PP_ALIGN.CENTER}])

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Cap-Track-PZPW.pptx")
prs.save(out)
print("OK ->", out, "| slajdów:", len(prs.slides._sldIdLst))
